#!/usr/bin/env python3
"""Generate all ElecSecure town hall presentation variants with branding."""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = Path(__file__).parent
LOGO_SMALL = BASE / "assets" / "elecsecure-logo-small.png"
LOGO_LARGE = BASE / "assets" / "elecsecure-logo.png"

NAVY = RGBColor(0x0B, 0x1F, 0x3A)
ELECTRIC_BLUE = RGBColor(0x00, 0x7A, 0xE6)
TEAL = RGBColor(0x00, 0xB8, 0xA9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF4, 0xF6, 0xF8)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
MID_GRAY = RGBColor(0x5A, 0x6A, 0x7A)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x42)

TITLE = {
    "layout": "title",
    "title": "ElecSecure",
    "subtitle": "Smart Electrical Safety & Energy Management\nAnnual Town Hall 2026",
    "footer": "Azfar Mushtaq | Managing Director",
    "notes": "Welcome everyone. Today I will introduce ElecSecure and explain why it matters to our organisation.",
}

CLOSING = {
    "layout": "closing",
    "title": "Thank You",
    "subtitle": "Questions & Discussion",
    "footer": "ElecSecure | Smart Electrical Safety & Energy Management",
    "notes": "Open floor for questions. Repeat each question for the room.",
}

CTA = {
    "layout": "cta",
    "title": "What We Need From You",
    "bullets": [
        "Leadership: endorse phased execution and resourcing",
        "Teams: identify collaboration and integration opportunities",
        "Everyone: connect us with potential customers and partners",
    ],
    "highlight": "Together, we make electricity safer and smarter.",
    "notes": "Clear call to action before Q&A.",
}

FULL_SLIDES = [
    TITLE,
    {"layout": "content", "title": "Today's Agenda", "bullets": [
        "Why electrical safety and efficiency matter now",
        "What ElecSecure is and how it works",
        "Market opportunity and competitive edge",
        "3-year roadmap and financial outlook",
        "Team, impact, and how you can support",
    ], "notes": "Set expectations: 20 minutes presentation, 10 minutes Q&A."},
    {"layout": "highlight", "title": "Why This Matters Now", "bullets": [
        "Electricity powers every home, workplace, and public facility",
        "Rising energy costs increase pressure to optimise usage",
        "Electrical faults remain a major fire and downtime risk",
        "Smart technology demand is accelerating across the UK",
    ], "highlight": "Safety  •  Savings  •  Control", "notes": "Lead with outcomes, not technology."},
    {"layout": "content", "title": "The Problem We Are Solving", "bullets": [
        "Limited real-time visibility into electrical health",
        "Faults are often detected too late",
        "Energy waste is hidden in daily operations",
        "Smart homes lack integrated electrical protection",
        "No unified platform for safety + efficiency + remote control",
    ], "notes": "Analogy: driving without dashboard lights."},
    {"layout": "two_column", "title": "Introducing ElecSecure", "left_title": "What it is", "left": [
        "Integrated smart electrical platform",
        "Hardware + mobile app + cloud analytics",
        "Built for homes and businesses",
    ], "right_title": "What it does", "right": [
        "Detects faults before they escalate",
        "Monitors energy use in real time",
        "Enables remote control and alerts",
        "Integrates with smart home systems",
    ], "notes": "Position as a platform, not just a device."},
    {"layout": "process_demo", "title": "How It Works", "steps": [
        ("1", "Sense", "Monitors voltage, current, temperature"),
        ("2", "Detect", "Identifies abnormal patterns"),
        ("3", "Alert", "Instant mobile notifications"),
        ("4", "Act", "Remote control & tripping"),
        ("5", "Optimise", "Insights reduce waste"),
    ], "demo_text": "▶  LIVE DEMO (90 sec): Dashboard → Alert → Remote Action",
    "notes": "KEY SLIDE. Optional live demo — see Live-Demo-Script.md. Walk through each step slowly."},
    {"layout": "cards", "title": "Who Benefits", "cards": [
        ("Families", "Safer homes\nLower bills"),
        ("Businesses", "Less downtime\nLower operating cost"),
        ("Electricians", "New service revenue\nTrusted product"),
        ("Public Sector", "Compliance\nSustainability goals"),
    ], "notes": "Even non-project staff can be connectors."},
    {"layout": "stats", "title": "Market Opportunity (UK)", "stats": [
        ("£2.4B+", "Market size 2023"),
        ("12.6%", "Projected annual growth"),
        ("2025", "UK smart meter rollout target"),
        ("High", "Demand for safety + efficiency"),
    ], "notes": "Keep numbers high-level."},
    {"layout": "comparison", "title": "Why ElecSecure Can Win", "headers": ["Typical alternatives", "ElecSecure"], "rows": [
        ["Thermostats / plugs only", "Full electrical system intelligence"],
        ["Passive protection", "Active monitoring + alerts"],
        ["Energy tracking only", "Safety + efficiency + control"],
        ["Disconnected devices", "Integrated smart ecosystem"],
    ], "notes": "We protect the system, not just measure it."},
    {"layout": "revenue", "title": "Business Model", "items": [
        ("Device Sales", "£300", "One-time hardware revenue"),
        ("Subscription", "£49/mo", "Cloud monitoring, updates, analytics"),
        ("Support Services", "£115", "Maintenance and technical assistance"),
    ], "notes": "Recurring revenue builds long-term value."},
    {"layout": "roadmap", "title": "3-Year Roadmap", "years": [
        ("Year 1", "Launch", "400 customers\nUK pilot regions\nCertification & brand build"),
        ("Year 2", "Scale", "800 customers\nNationwide growth\nProfitability"),
        ("Year 3", "Lead", "1,500 customers\nInternational expansion\nMarket leadership"),
    ], "notes": "KEY SLIDE. Pause after each year for emphasis."},
    {"layout": "financial", "title": "Financial Outlook (High Level)", "rows": [
        ("Year 1", "£235K revenue", "Foundation year"),
        ("Year 2", "£751K revenue", "£154K profit after tax"),
        ("Year 3", "£1.42M revenue", "£448K profit after tax"),
    ], "notes": "KEY SLIDE. Year 1 is investment by design. Profitability from Year 2."},
    {"layout": "team", "title": "Leadership & Delivery Team", "bullets": [
        "Managing Director: Azfar Mushtaq",
        "Software development and mobile platform",
        "Electrical engineering and QA",
        "Data science and analytics",
        "Sales, marketing, and customer operations",
    ], "notes": "Cross-functional execution capability."},
    {"layout": "impact", "title": "Impact Beyond Revenue", "bullets": [
        "Job creation: installers, support, engineering, sales",
        "Reduced electrical fire and downtime risk",
        "Lower energy waste and carbon footprint",
        "Supports UK safety and efficiency policy goals",
    ], "notes": "End on mission and social value."},
    {"layout": "risks", "title": "Risks & Mitigation", "pairs": [
        ("Low brand awareness", "Targeted launch + partner channels"),
        ("Upfront investment", "Phased rollout + strategic funding"),
        ("Regulatory complexity", "Early UK certification pathway"),
        ("Fast tech change", "Continuous R&D and product updates"),
    ], "notes": "Naming risks builds credibility."},
    CTA,
    CLOSING,
]

SHORT_10_SLIDES = [
    TITLE,
    {"layout": "highlight", "title": "Why Now & The Problem", "bullets": [
        "Electricity powers everything — but we lack real-time visibility",
        "Faults are detected too late; energy waste stays hidden",
        "Rising costs and safety risks demand a smarter approach",
        "Smart home tech rarely includes true electrical protection",
    ], "highlight": "Safety  •  Savings  •  Control", "notes": "Combined problem + urgency slide. 2 minutes."},
    {"layout": "two_column", "title": "ElecSecure Solution", "left_title": "What it is", "left": [
        "Smart device + mobile app + cloud platform",
        "For homes and businesses",
        "Real-time safety and efficiency",
    ], "right_title": "What it does", "right": [
        "Detects faults before they escalate",
        "Sends instant alerts to your phone",
        "Monitors and controls power remotely",
        "Integrates with smart home systems",
    ], "notes": "3 minutes. Core product explanation."},
    {"layout": "process_demo", "title": "How It Works", "steps": [
        ("1", "Sense", "Monitor electrical health"),
        ("2", "Detect", "Spot abnormal patterns"),
        ("3", "Alert", "Notify user instantly"),
        ("4", "Act", "Remote control & protection"),
        ("5", "Optimise", "Reduce waste & cost"),
    ], "demo_text": "▶  LIVE DEMO (90 sec) — see demo script",
    "notes": "KEY SLIDE. 3 minutes. Demo or walkthrough."},
    {"layout": "cards", "title": "Who Benefits", "cards": [
        ("Families", "Safer homes\nLower bills"),
        ("Businesses", "Less downtime\nLower cost"),
        ("Electricians", "New revenue\nTrusted product"),
        ("Public Sector", "Compliance\nSustainability"),
    ], "notes": "2 minutes. Inclusive messaging."},
    {"layout": "stats", "title": "Market & Competitive Edge", "stats": [
        ("£2.4B+", "UK market size"),
        ("12.6%", "Annual growth"),
        ("3-in-1", "Safety + efficiency + control"),
        ("First", "Integrated platform advantage"),
    ], "notes": "2 minutes. Market validation."},
    {"layout": "revenue", "title": "Business Model", "items": [
        ("Device Sales", "£300", "Hardware revenue"),
        ("Subscription", "£49/mo", "Recurring cloud services"),
        ("Support", "£115", "Maintenance packages"),
    ], "notes": "1.5 minutes."},
    {"layout": "roadmap_financial", "title": "Roadmap & Financial Outlook", "years": [
        ("Year 1", "400 customers", "£235K revenue"),
        ("Year 2", "800 customers", "£751K | £154K profit"),
        ("Year 3", "1,500 customers", "£1.42M | £448K profit"),
    ], "notes": "KEY SLIDE. 3 minutes. Pause on numbers."},
    {"layout": "risks", "title": "Risks & How We Manage Them", "pairs": [
        ("Low brand awareness", "Partner-led launch"),
        ("Upfront investment", "Phased rollout"),
        ("Regulation", "Early certification"),
        ("Tech change", "Continuous R&D"),
    ], "notes": "1.5 minutes."},
    {"layout": "cta", "title": "What We Need From You", "bullets": [
        "Leadership: endorse phased execution",
        "Teams: identify collaboration opportunities",
        "Everyone: connect us with customers and partners",
    ], "highlight": "Together, we make electricity safer and smarter.",
    "notes": "1 minute. Then Q&A."},
]

EXECUTIVE_5_SLIDES = [
    {**TITLE, "subtitle": "Executive Briefing | Annual Town Hall 2026", "notes": "5-minute executive version. Be concise and confident."},
    {"layout": "highlight", "title": "The Opportunity", "bullets": [
        "UK electrical safety & efficiency market: £2.4B+, growing 12.6% annually",
        "Current systems lack real-time safety visibility and energy optimisation",
        "ElecSecure: smart device + app + cloud — Safety, Savings, Control",
        "Targets homes, businesses, and public sector across the UK",
    ], "highlight": "Problem → Platform → Market fit", "notes": "60 seconds. Problem and solution together."},
    {"layout": "revenue", "title": "Business Model & Differentiation", "items": [
        ("Device", "£300", "Hardware sale per installation"),
        ("Subscription", "£49/mo", "Recurring cloud & analytics revenue"),
        ("Support", "£115", "Maintenance and service packages"),
    ], "notes": "60 seconds. Emphasise recurring revenue model."},
    {"layout": "roadmap_financial", "title": "3-Year Plan & Returns", "years": [
        ("Year 1", "400 customers", "£235K — Foundation"),
        ("Year 2", "800 customers", "£751K — Profitable"),
        ("Year 3", "1,500 customers", "£1.42M — Market leader"),
    ], "notes": "90 seconds. This is the money slide for executives."},
    {"layout": "cta", "title": "Decision Requested", "bullets": [
        "Endorse phased UK launch and resourcing plan",
        "Enable cross-functional support across teams",
        "Champion introductions to strategic partners and customers",
    ], "highlight": "Approve direction. Accelerate execution.",
    "notes": "60 seconds. Clear ask. Offer full deck for detail."},
]

VARIANTS = {
    "full": ("ElecSecure-Townhall-2026-Full.pptx", FULL_SLIDES),
    "short": ("ElecSecure-Townhall-2026-10-Slides.pptx", SHORT_10_SLIDES),
    "executive": ("ElecSecure-Townhall-2026-Executive-5min.pptx", EXECUTIVE_5_SLIDES),
}


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_logo(slide, large=False):
    logo = LOGO_LARGE if large and LOGO_LARGE.exists() else LOGO_SMALL
    if not logo.exists():
        return
    if large:
        slide.shapes.add_picture(str(logo), Inches(5.4), Inches(0.55), height=Inches(1.3))
    else:
        slide.shapes.add_picture(str(logo), Inches(12.0), Inches(0.18), height=Inches(0.75))


def add_header_bar(slide, prs, title_text):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.1), prs.slide_width, Inches(0.08))
    accent.fill.solid(); accent.fill.fore_color.rgb = ELECTRIC_BLUE; accent.line.fill.background()
    title = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(11.2), Inches(0.7))
    tf = title.text_frame; tf.text = title_text
    p = tf.paragraphs[0]; p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Calibri"
    add_logo(slide)


def add_bullets(slide, bullets, top=1.6, left=0.8, width=11.5, size=22, color=DARK_TEXT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5))
    tf = box.text_frame; tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet; p.level = 0; p.font.size = Pt(size); p.font.color.rgb = color
        p.font.name = "Calibri"; p.space_after = Pt(14); p.bullet = True


def add_notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text


def build_title_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(slide, NAVY)
    add_logo(slide, large=True)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(4.8), prs.slide_width, Inches(0.12))
    accent.fill.solid(); accent.fill.fore_color.rgb = TEAL; accent.line.fill.background()
    title = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.2))
    tf = title.text_frame; tf.text = data["title"]
    p = tf.paragraphs[0]; p.font.size = Pt(54); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = "Calibri"
    sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11.5), Inches(1.5))
    tf = sub.text_frame; tf.text = data["subtitle"]
    for p in tf.paragraphs: p.font.size = Pt(24); p.font.color.rgb = RGBColor(0xCC, 0xE5, 0xFF); p.font.name = "Calibri"
    if data.get("footer"):
        foot = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11), Inches(0.5))
        tf = foot.text_frame; tf.text = data["footer"]
        p = tf.paragraphs[0]; p.font.size = Pt(16); p.font.color.rgb = MID_GRAY; p.font.name = "Calibri"
    add_notes(slide, data.get("notes", ""))


def build_content_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(slide, WHITE)
    add_header_bar(slide, prs, data["title"]); add_bullets(slide, data["bullets"]); add_notes(slide, data.get("notes", ""))


def build_highlight_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(slide, WHITE)
    add_header_bar(slide, prs, data["title"]); add_bullets(slide, data["bullets"], top=1.5)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.0), Inches(10.3), Inches(1.0))
    box.fill.solid(); box.fill.fore_color.rgb = ELECTRIC_BLUE; box.line.fill.background()
    tf = box.text_frame; tf.text = data["highlight"]; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.font.size = Pt(28); p.font.bold = True
    p.font.color.rgb = WHITE; p.font.name = "Calibri"
    add_notes(slide, data.get("notes", ""))


def build_two_column_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(slide, WHITE)
    add_header_bar(slide, prs, data["title"])
    for title, items, left in [(data["left_title"], data["left"], 0.7), (data["right_title"], data["right"], 6.5)]:
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.5), Inches(5.8), Inches(4.8))
        card.fill.solid(); card.fill.fore_color.rgb = LIGHT_GRAY; card.line.color.rgb = RGBColor(0xDD, 0xE3, 0xEA)
        tbox = slide.shapes.add_textbox(Inches(left + 0.3), Inches(1.8), Inches(5.2), Inches(4.2))
        tf = tbox.text_frame; tf.word_wrap = True
        h = tf.paragraphs[0]; h.text = title; h.font.size = Pt(22); h.font.bold = True; h.font.color.rgb = NAVY; h.space_after = Pt(12)
        for item in items:
            p = tf.add_paragraph(); p.text = item; p.font.size = Pt(18); p.font.color.rgb = DARK_TEXT; p.bullet = True; p.space_after = Pt(8)
    add_notes(slide, data.get("notes", ""))


def build_process_demo_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(slide, WHITE)
    add_header_bar(slide, prs, data["title"])
    x_positions = [0.5, 2.7, 4.9, 7.1, 9.3]
    for i, (num, label, desc) in enumerate(data["steps"]):
        left = x_positions[i]
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(1.9), Inches(0.9), Inches(0.9))
        circle.fill.solid(); circle.fill.fore_color.rgb = ELECTRIC_BLUE; circle.line.fill.background()
        tf = circle.text_frame; tf.text = num; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = WHITE
        label_box = slide.shapes.add_textbox(Inches(left - 0.2), Inches(2.9), Inches(2.0), Inches(0.5))
        tf = label_box.text_frame; tf.text = label
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = NAVY
        desc_box = slide.shapes.add_textbox(Inches(left - 0.35), Inches(3.35), Inches(2.3), Inches(1.2))
        tf = desc_box.text_frame; tf.word_wrap = True; tf.text = desc
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.font.size = Pt(12); p.font.color.rgb = MID_GRAY
        if i < len(data["steps"]) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(left + 1.0), Inches(2.25), Inches(0.7), Inches(0.3))
            arrow.fill.solid(); arrow.fill.fore_color.rgb = TEAL; arrow.line.fill.background()
    demo = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(5.15), Inches(10.9), Inches(0.85))
    demo.fill.solid(); demo.fill.fore_color.rgb = ACCENT_ORANGE; demo.line.fill.background()
    tf = demo.text_frame; tf.text = data.get("demo_text", "▶  LIVE DEMO"); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = WHITE
    add_notes(slide, data.get("notes", ""))


def build_cards_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(slide, WHITE)
    add_header_bar(slide, prs, data["title"])
    positions = [(0.7, 1.6), (6.5, 1.6), (0.7, 4.2), (6.5, 4.2)]
    colors = [NAVY, ELECTRIC_BLUE, TEAL, ACCENT_ORANGE]
    for i, (title, body) in enumerate(data["cards"]):
        left, top = positions[i]
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(5.8), Inches(2.2))
        card.fill.solid(); card.fill.fore_color.rgb = colors[i]; card.line.fill.background()
        tbox = slide.shapes.add_textbox(Inches(left + 0.3), Inches(top + 0.3), Inches(5.2), Inches(1.6))
        tf = tbox.text_frame; tf.word_wrap = True
        h = tf.paragraphs[0]; h.text = title; h.font.size = Pt(22); h.font.bold = True; h.font.color.rgb = WHITE; h.space_after = Pt(8)
        for line in body.split("\n"):
            p = tf.add_paragraph(); p.text = line; p.font.size = Pt(16); p.font.color.rgb = WHITE
    add_notes(slide, data.get("notes", ""))


def build_stats_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(slide, WHITE)
    add_header_bar(slide, prs, data["title"])
    positions = [(0.8, 1.8), (6.6, 1.8), (0.8, 4.3), (6.6, 4.3)]
    for (value, label), (left, top) in zip(data["stats"], positions):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(5.5), Inches(2.0))
        card.fill.solid(); card.fill.fore_color.rgb = LIGHT_GRAY; card.line.color.rgb = RGBColor(0xDD, 0xE3, 0xEA)
        vbox = slide.shapes.add_textbox(Inches(left), Inches(top + 0.35), Inches(5.5), Inches(0.8))
        tf = vbox.text_frame; tf.text = value
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = ELECTRIC_BLUE
        lbox = slide.shapes.add_textbox(Inches(left), Inches(top + 1.2), Inches(5.5), Inches(0.6))
        tf = lbox.text_frame; tf.text = label
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.font.size = Pt(16); p.font.color.rgb = DARK_TEXT
    add_notes(slide, data.get("notes", ""))


def build_comparison_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(slide, WHITE)
    add_header_bar(slide, prs, data["title"])
    table = slide.shapes.add_table(5, 2, Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.5)).table
    table.columns[0].width = Inches(5.8); table.columns[1].width = Inches(5.9)
    for col, header in enumerate(data["headers"]):
        cell = table.cell(0, col); cell.text = header
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY if col == 1 else MID_GRAY
        p = cell.text_frame.paragraphs[0]; p.font.bold = True; p.font.size = Pt(16); p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    for r, row in enumerate(data["rows"], start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c); cell.text = val
            if c == 1: cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFD)
            p = cell.text_frame.paragraphs[0]; p.font.size = Pt(14); p.font.color.rgb = DARK_TEXT; p.alignment = PP_ALIGN.CENTER
    add_notes(slide, data.get("notes", ""))


def build_revenue_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(slide, WHITE)
    add_header_bar(slide, prs, data["title"])
    for i, (title, price, desc) in enumerate(data["items"]):
        top = 1.7 + i * 1.7
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(top), Inches(11.3), Inches(1.4))
        card.fill.solid(); card.fill.fore_color.rgb = LIGHT_GRAY; card.line.color.rgb = RGBColor(0xDD, 0xE3, 0xEA)
        for box_left, text, size, bold, color in [
            (1.3, title, 20, True, NAVY), (5.8, price, 26, True, ELECTRIC_BLUE), (8.0, desc, 14, False, MID_GRAY)
        ]:
            b = slide.shapes.add_textbox(Inches(box_left), Inches(top + 0.2), Inches(4.5), Inches(1.0))
            tf = b.text_frame; tf.text = text
            p = tf.paragraphs[0]; p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color
    add_notes(slide, data.get("notes", ""))


def build_roadmap_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(slide, WHITE)
    add_header_bar(slide, prs, data["title"])
    colors = [ELECTRIC_BLUE, TEAL, NAVY]
    for i, (year, phase, details) in enumerate(data["years"]):
        left = 0.7 + i * 4.1
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.8), Inches(3.7), Inches(4.5))
        card.fill.solid(); card.fill.fore_color.rgb = colors[i]; card.line.fill.background()
        tbox = slide.shapes.add_textbox(Inches(left + 0.2), Inches(2.0), Inches(3.3), Inches(4.0))
        tf = tbox.text_frame; tf.word_wrap = True
        for j, line in enumerate([year, phase, ""] + details.split("\n")):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph(); p.text = line
            p.font.size = Pt(24 if j == 0 else 18 if j == 1 else 14); p.font.bold = j < 2
            p.font.color.rgb = WHITE; p.bullet = bool(line) and j > 2
    add_notes(slide, data.get("notes", ""))


def build_roadmap_financial_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(slide, WHITE)
    add_header_bar(slide, prs, data["title"])
    colors = [ELECTRIC_BLUE, TEAL, NAVY]
    for i, (year, customers, financial) in enumerate(data["years"]):
        left = 0.7 + i * 4.1
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.8), Inches(3.7), Inches(4.5))
        card.fill.solid(); card.fill.fore_color.rgb = colors[i]; card.line.fill.background()
        tbox = slide.shapes.add_textbox(Inches(left + 0.2), Inches(2.2), Inches(3.3), Inches(3.8))
        tf = tbox.text_frame; tf.word_wrap = True
        for j, line in enumerate([year, customers, financial]):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph(); p.text = line
            p.font.size = Pt(26 if j == 0 else 18 if j == 1 else 16); p.font.bold = j < 2
            p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    add_notes(slide, data.get("notes", ""))


def build_financial_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(slide, WHITE)
    add_header_bar(slide, prs, data["title"])
    for i, (year, revenue, note) in enumerate(data["rows"]):
        top = 1.8 + i * 1.5
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(top), Inches(11.7), Inches(1.2))
        bar.fill.solid(); bar.fill.fore_color.rgb = LIGHT_GRAY; bar.line.color.rgb = RGBColor(0xDD, 0xE3, 0xEA)
        for left, text, size, color, bold in [(1.1, year, 20, NAVY, True), (2.8, revenue, 20, ELECTRIC_BLUE, True), (7.2, note, 14, MID_GRAY, False)]:
            b = slide.shapes.add_textbox(Inches(left), Inches(top + 0.25), Inches(4.8), Inches(0.7))
            tf = b.text_frame; tf.text = text
            p = tf.paragraphs[0]; p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color
    add_notes(slide, data.get("notes", ""))


def build_risks_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(slide, WHITE)
    add_header_bar(slide, prs, data["title"])
    for i, (risk, mitigation) in enumerate(data["pairs"]):
        top = 1.6 + i * 1.25
        rbox = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(5.2), Inches(0.9))
        tf = rbox.text_frame; tf.text = f"Risk: {risk}"
        p = tf.paragraphs[0]; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = ACCENT_ORANGE
        mbox = slide.shapes.add_textbox(Inches(6.2), Inches(top), Inches(6.0), Inches(0.9))
        tf = mbox.text_frame; tf.text = f"Plan: {mitigation}"
        p = tf.paragraphs[0]; p.font.size = Pt(16); p.font.color.rgb = TEAL
    add_notes(slide, data.get("notes", ""))


def build_cta_slide(prs, data):
    build_highlight_slide(prs, data)


def build_closing_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); set_slide_bg(slide, NAVY)
    add_logo(slide, large=True)
    title = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.5), Inches(1.2))
    tf = title.text_frame; tf.text = data["title"]
    p = tf.paragraphs[0]; p.font.size = Pt(48); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(11.5), Inches(1.0))
    tf = sub.text_frame; tf.text = data["subtitle"]
    p = tf.paragraphs[0]; p.font.size = Pt(28); p.font.color.rgb = TEAL; p.alignment = PP_ALIGN.CENTER
    if data.get("footer"):
        foot = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5))
        tf = foot.text_frame; tf.text = data["footer"]
        p = tf.paragraphs[0]; p.font.size = Pt(14); p.font.color.rgb = MID_GRAY; p.alignment = PP_ALIGN.CENTER
    add_notes(slide, data.get("notes", ""))


BUILDERS = {
    "title": build_title_slide, "content": build_content_slide, "highlight": build_highlight_slide,
    "two_column": build_two_column_slide, "process_demo": build_process_demo_slide,
    "cards": build_cards_slide, "stats": build_stats_slide, "comparison": build_comparison_slide,
    "revenue": build_revenue_slide, "roadmap": build_roadmap_slide, "roadmap_financial": build_roadmap_financial_slide,
    "financial": build_financial_slide, "team": build_content_slide, "impact": build_content_slide,
    "risks": build_risks_slide, "cta": build_cta_slide, "closing": build_closing_slide,
}


def generate(slides, output_path):
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    for slide_data in slides:
        BUILDERS[slide_data["layout"]](prs, slide_data)
    prs.save(str(output_path))
    return len(slides)


def main():
    parser = argparse.ArgumentParser(description="Generate ElecSecure presentation variants")
    parser.add_argument("--variant", choices=["all", "full", "short", "executive"], default="all")
    args = parser.parse_args()

    # Ensure logo exists
    import subprocess
    subprocess.run(["python3", str(BASE / "assets" / "create_logo.py")], check=True)

    targets = VARIANTS.items() if args.variant == "all" else [(args.variant, VARIANTS[args.variant])]
    for key, (filename, slides) in targets:
        path = BASE / filename
        count = generate(slides, path)
        print(f"Created: {path} ({count} slides)")

    # Also create legacy filename alias
    if args.variant in ("all", "full"):
        import shutil
        shutil.copy(BASE / "ElecSecure-Townhall-2026-Full.pptx", BASE / "ElecSecure-Townhall-2026.pptx")
        print(f"Created: {BASE / 'ElecSecure-Townhall-2026.pptx'} (alias)")


if __name__ == "__main__":
    main()
