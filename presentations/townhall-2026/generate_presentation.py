#!/usr/bin/env python3
"""Generate ElecSecure Town Hall PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colours
NAVY = RGBColor(0x0B, 0x1F, 0x3A)
ELECTRIC_BLUE = RGBColor(0x00, 0x7A, 0xE6)
TEAL = RGBColor(0x00, 0xB8, 0xA9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF4, 0xF6, 0xF8)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
MID_GRAY = RGBColor(0x5A, 0x6A, 0x7A)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x42)

OUTPUT = "/workspace/presentations/townhall-2026/ElecSecure-Townhall-2026.pptx"

SLIDES = [
    {
        "layout": "title",
        "title": "ElecSecure",
        "subtitle": "Smart Electrical Safety & Energy Management\nAnnual Town Hall 2026",
        "footer": "Azfar Mushtaq | Managing Director",
        "notes": "Welcome everyone. Today I will introduce ElecSecure and explain why it matters to our organisation and communities.",
    },
    {
        "layout": "content",
        "title": "Today's Agenda",
        "bullets": [
            "Why electrical safety and efficiency matter now",
            "What ElecSecure is and how it works",
            "Market opportunity and competitive edge",
            "3-year roadmap and financial outlook",
            "Team, impact, and how you can support",
        ],
        "notes": "Set expectations: 20 minutes presentation, 10 minutes Q&A.",
    },
    {
        "layout": "highlight",
        "title": "Why This Matters Now",
        "bullets": [
            "Electricity powers every home, workplace, and public facility",
            "Rising energy costs increase pressure to optimise usage",
            "Electrical faults remain a major fire and downtime risk",
            "Smart technology demand is accelerating across the UK",
        ],
        "highlight": "Safety  •  Savings  •  Control",
        "notes": "Lead with outcomes, not technology. Repeat the phrase: Safety, Savings, Control.",
    },
    {
        "layout": "content",
        "title": "The Problem We Are Solving",
        "bullets": [
            "Limited real-time visibility into electrical health",
            "Faults are often detected too late",
            "Energy waste is hidden in daily operations",
            "Smart homes lack integrated electrical protection",
            "No unified platform for safety + efficiency + remote control",
        ],
        "notes": "Use analogy: driving without dashboard lights.",
    },
    {
        "layout": "two_column",
        "title": "Introducing ElecSecure",
        "left_title": "What it is",
        "left": [
            "Integrated smart electrical platform",
            "Hardware + mobile app + cloud analytics",
            "Built for homes and businesses",
        ],
        "right_title": "What it does",
        "right": [
            "Detects faults before they escalate",
            "Monitors energy use in real time",
            "Enables remote control and alerts",
            "Integrates with smart home systems",
        ],
        "notes": "Position as a platform, not just a device.",
    },
    {
        "layout": "process",
        "title": "How It Works",
        "steps": [
            ("1", "Sense", "Device monitors voltage, current, temperature"),
            ("2", "Detect", "AI + rules identify abnormal patterns"),
            ("3", "Alert", "Users receive instant mobile notifications"),
            ("4", "Act", "Remote control and protective tripping"),
            ("5", "Optimise", "Insights reduce waste and cost"),
        ],
        "notes": "Walk through each step slowly. This slide is key for non-technical audience.",
    },
    {
        "layout": "cards",
        "title": "Who Benefits",
        "cards": [
            ("Families", "Safer homes\nLower bills"),
            ("Businesses", "Less downtime\nLower operating cost"),
            ("Electricians", "New service revenue\nTrusted product"),
            ("Public Sector", "Compliance\nSustainability goals"),
        ],
        "notes": "Make it inclusive: even non-project staff can be connectors.",
    },
    {
        "layout": "stats",
        "title": "Market Opportunity (UK)",
        "stats": [
            ("£2.4B+", "Market size 2023"),
            ("12.6%", "Projected annual growth"),
            ("2025", "UK smart meter rollout target"),
            ("High", "Demand for safety + efficiency"),
        ],
        "notes": "Keep numbers high-level. Emphasise timing advantage.",
    },
    {
        "layout": "comparison",
        "title": "Why ElecSecure Can Win",
        "headers": ["Typical alternatives", "ElecSecure"],
        "rows": [
            ["Thermostats / plugs only", "Full electrical system intelligence"],
            ["Passive protection", "Active monitoring + alerts"],
            ["Energy tracking only", "Safety + efficiency + control"],
            ["Disconnected devices", "Integrated smart ecosystem"],
        ],
        "notes": "One-liner: We protect the system, not just measure it.",
    },
    {
        "layout": "revenue",
        "title": "Business Model",
        "items": [
            ("Device Sales", "£300", "One-time hardware revenue"),
            ("Subscription", "£49/mo", "Cloud monitoring, updates, analytics"),
            ("Support Services", "£115", "Maintenance and technical assistance"),
        ],
        "notes": "Explain recurring revenue builds long-term value.",
    },
    {
        "layout": "roadmap",
        "title": "3-Year Roadmap",
        "years": [
            ("Year 1", "Launch", "400 customers\nUK pilot regions\nCertification & brand build"),
            ("Year 2", "Scale", "800 customers\nNationwide growth\nProfitability"),
            ("Year 3", "Lead", "1,500 customers\nInternational expansion\nMarket leadership"),
        ],
        "notes": "Pause after each year for emphasis.",
    },
    {
        "layout": "financial",
        "title": "Financial Outlook (High Level)",
        "rows": [
            ("Year 1", "£235K revenue", "Foundation year"),
            ("Year 2", "£751K revenue", "£154K profit after tax"),
            ("Year 3", "£1.42M revenue", "£448K profit after tax"),
        ],
        "notes": "Year 1 is investment by design. Profitability from Year 2.",
    },
    {
        "layout": "team",
        "title": "Leadership & Delivery Team",
        "bullets": [
            "Managing Director: Azfar Mushtaq",
            "Software development and mobile platform",
            "Electrical engineering and QA",
            "Data science and analytics",
            "Sales, marketing, and customer operations",
        ],
        "notes": "Emphasise cross-functional execution capability.",
    },
    {
        "layout": "impact",
        "title": "Impact Beyond Revenue",
        "bullets": [
            "Job creation: installers, support, engineering, sales",
            "Reduced electrical fire and downtime risk",
            "Lower energy waste and carbon footprint",
            "Supports UK safety and efficiency policy goals",
        ],
        "notes": "End this section on mission and social value.",
    },
    {
        "layout": "risks",
        "title": "Risks & Mitigation",
        "pairs": [
            ("Low brand awareness", "Targeted launch + partner channels"),
            ("Upfront investment", "Phased rollout + strategic funding"),
            ("Regulatory complexity", "Early UK certification pathway"),
            ("Fast tech change", "Continuous R&D and product updates"),
        ],
        "notes": "Naming risks builds leadership credibility.",
    },
    {
        "layout": "cta",
        "title": "What We Need From You",
        "bullets": [
            "Leadership: endorse phased execution and resourcing",
            "Teams: identify collaboration and integration opportunities",
            "Everyone: connect us with potential customers and partners",
        ],
        "highlight": "Together, we make electricity safer and smarter.",
        "notes": "Clear call to action before Q&A.",
    },
    {
        "layout": "closing",
        "title": "Thank You",
        "subtitle": "Questions & Discussion",
        "footer": "ElecSecure | Smart Electrical Safety & Energy Management",
        "notes": "Open floor for questions. Repeat each question for the room.",
    },
]


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_header_bar(slide, prs, title_text):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.1)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.1), prs.slide_width, Inches(0.08)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = ELECTRIC_BLUE
    accent.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(12), Inches(0.7))
    tf = title.text_frame
    tf.text = title_text
    p = tf.paragraphs[0]
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"


def add_bullets(slide, bullets, top=1.6, left=0.8, width=11.5, size=22, color=DARK_TEXT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(14)
        p.bullet = True


def add_notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text


def build_title_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(4.8), prs.slide_width, Inches(0.12)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(1.2))
    tf = title.text_frame
    tf.text = data["title"]
    p = tf.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(11.5), Inches(1.5))
    tf = sub.text_frame
    tf.text = data["subtitle"]
    for p in tf.paragraphs:
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(0xCC, 0xE5, 0xFF)
        p.font.name = "Calibri"

    if data.get("footer"):
        foot = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11), Inches(0.5))
        tf = foot.text_frame
        tf.text = data["footer"]
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = MID_GRAY
        p.font.name = "Calibri"

    add_notes(slide, data.get("notes", ""))


def build_content_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header_bar(slide, prs, data["title"])
    add_bullets(slide, data["bullets"])
    add_notes(slide, data.get("notes", ""))


def build_highlight_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header_bar(slide, prs, data["title"])
    add_bullets(slide, data["bullets"], top=1.5)

    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.0), Inches(10.3), Inches(1.0)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = ELECTRIC_BLUE
    box.line.fill.background()
    tf = box.text_frame
    tf.text = data["highlight"]
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"

    add_notes(slide, data.get("notes", ""))


def build_two_column_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header_bar(slide, prs, data["title"])

    for idx, (title, items, left) in enumerate([
        (data["left_title"], data["left"], 0.7),
        (data["right_title"], data["right"], 6.5),
    ]):
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.5), Inches(5.8), Inches(4.8)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = RGBColor(0xDD, 0xE3, 0xEA)

        tbox = slide.shapes.add_textbox(Inches(left + 0.3), Inches(1.8), Inches(5.2), Inches(4.2))
        tf = tbox.text_frame
        tf.word_wrap = True
        h = tf.paragraphs[0]
        h.text = title
        h.font.size = Pt(22)
        h.font.bold = True
        h.font.color.rgb = NAVY
        h.font.name = "Calibri"
        h.space_after = Pt(12)

        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = DARK_TEXT
            p.font.name = "Calibri"
            p.level = 0
            p.bullet = True
            p.space_after = Pt(8)

    add_notes(slide, data.get("notes", ""))


def build_process_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header_bar(slide, prs, data["title"])

    x_positions = [0.5, 2.7, 4.9, 7.1, 9.3]
    for i, (num, label, desc) in enumerate(data["steps"]):
        left = x_positions[i]
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(left), Inches(2.0), Inches(0.9), Inches(0.9)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = ELECTRIC_BLUE
        circle.line.fill.background()
        tf = circle.text_frame
        tf.text = num
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE

        label_box = slide.shapes.add_textbox(Inches(left - 0.2), Inches(3.1), Inches(2.0), Inches(0.5))
        tf = label_box.text_frame
        tf.text = label
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = NAVY

        desc_box = slide.shapes.add_textbox(Inches(left - 0.35), Inches(3.6), Inches(2.3), Inches(2.0))
        tf = desc_box.text_frame
        tf.word_wrap = True
        tf.text = desc
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(13)
        p.font.color.rgb = MID_GRAY

        if i < len(data["steps"]) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, Inches(left + 1.0), Inches(2.35), Inches(0.7), Inches(0.3)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = TEAL
            arrow.line.fill.background()

    add_notes(slide, data.get("notes", ""))


def build_cards_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header_bar(slide, prs, data["title"])

    positions = [(0.7, 1.6), (6.5, 1.6), (0.7, 4.2), (6.5, 4.2)]
    colors = [NAVY, ELECTRIC_BLUE, TEAL, ACCENT_ORANGE]

    for i, (title, body) in enumerate(data["cards"]):
        left, top = positions[i]
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(5.8), Inches(2.2)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = colors[i]
        card.line.fill.background()

        tbox = slide.shapes.add_textbox(Inches(left + 0.3), Inches(top + 0.3), Inches(5.2), Inches(1.6))
        tf = tbox.text_frame
        tf.word_wrap = True
        h = tf.paragraphs[0]
        h.text = title
        h.font.size = Pt(22)
        h.font.bold = True
        h.font.color.rgb = WHITE
        h.space_after = Pt(8)

        for line in body.split("\n"):
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(16)
            p.font.color.rgb = WHITE
            p.font.name = "Calibri"

    add_notes(slide, data.get("notes", ""))


def build_stats_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header_bar(slide, prs, data["title"])

    positions = [(0.8, 1.8), (6.6, 1.8), (0.8, 4.3), (6.6, 4.3)]
    for i, (value, label) in enumerate(data["stats"]):
        left, top = positions[i]
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(5.5), Inches(2.0)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = RGBColor(0xDD, 0xE3, 0xEA)

        vbox = slide.shapes.add_textbox(Inches(left), Inches(top + 0.35), Inches(5.5), Inches(0.8))
        tf = vbox.text_frame
        tf.text = value
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = ELECTRIC_BLUE

        lbox = slide.shapes.add_textbox(Inches(left), Inches(top + 1.2), Inches(5.5), Inches(0.6))
        tf = lbox.text_frame
        tf.text = label
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_TEXT

    add_notes(slide, data.get("notes", ""))


def build_comparison_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header_bar(slide, prs, data["title"])

    table = slide.shapes.add_table(5, 2, Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.5)).table
    table.columns[0].width = Inches(5.8)
    table.columns[1].width = Inches(5.9)

    for col, header in enumerate(data["headers"]):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY if col == 1 else MID_GRAY
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    for r, row in enumerate(data["rows"], start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            if c == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFD)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_TEXT
            p.alignment = PP_ALIGN.CENTER

    add_notes(slide, data.get("notes", ""))


def build_revenue_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header_bar(slide, prs, data["title"])

    for i, (title, price, desc) in enumerate(data["items"]):
        top = 1.7 + i * 1.7
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(top), Inches(11.3), Inches(1.4)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = RGBColor(0xDD, 0xE3, 0xEA)

        tbox = slide.shapes.add_textbox(Inches(1.3), Inches(top + 0.2), Inches(4.5), Inches(1.0))
        tf = tbox.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = NAVY

        pbox = slide.shapes.add_textbox(Inches(5.8), Inches(top + 0.15), Inches(2.0), Inches(1.0))
        tf = pbox.text_frame
        tf.text = price
        p = tf.paragraphs[0]
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = ELECTRIC_BLUE

        dbox = slide.shapes.add_textbox(Inches(8.0), Inches(top + 0.25), Inches(4.0), Inches(1.0))
        tf = dbox.text_frame
        tf.text = desc
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = MID_GRAY

    add_notes(slide, data.get("notes", ""))


def build_roadmap_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header_bar(slide, prs, data["title"])

    colors = [ELECTRIC_BLUE, TEAL, NAVY]
    for i, (year, phase, details) in enumerate(data["years"]):
        left = 0.7 + i * 4.1
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.8), Inches(3.7), Inches(4.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = colors[i]
        card.line.fill.background()

        tbox = slide.shapes.add_textbox(Inches(left + 0.2), Inches(2.0), Inches(3.3), Inches(4.0))
        tf = tbox.text_frame
        tf.word_wrap = True
        lines = [year, phase, ""] + details.split("\n")
        for j, line in enumerate(lines):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = line
            if j == 0:
                p.font.size = Pt(24)
                p.font.bold = True
            elif j == 1:
                p.font.size = Pt(18)
                p.font.bold = True
            else:
                p.font.size = Pt(14)
                p.bullet = bool(line)
            p.font.color.rgb = WHITE
            p.space_after = Pt(6)

    add_notes(slide, data.get("notes", ""))


def build_financial_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header_bar(slide, prs, data["title"])

    for i, (year, revenue, note) in enumerate(data["rows"]):
        top = 1.8 + i * 1.5
        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(top), Inches(11.7), Inches(1.2)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = LIGHT_GRAY
        bar.line.color.rgb = RGBColor(0xDD, 0xE3, 0xEA)

        ybox = slide.shapes.add_textbox(Inches(1.1), Inches(top + 0.25), Inches(1.5), Inches(0.7))
        tf = ybox.text_frame
        tf.text = year
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = NAVY

        rbox = slide.shapes.add_textbox(Inches(2.8), Inches(top + 0.25), Inches(4.0), Inches(0.7))
        tf = rbox.text_frame
        tf.text = revenue
        p = tf.paragraphs[0]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = ELECTRIC_BLUE

        nbox = slide.shapes.add_textbox(Inches(7.2), Inches(top + 0.3), Inches(4.8), Inches(0.7))
        tf = nbox.text_frame
        tf.text = note
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = MID_GRAY

    add_notes(slide, data.get("notes", ""))


def build_team_slide(prs, data):
    build_content_slide(prs, data)


def build_impact_slide(prs, data):
    build_content_slide(prs, data)


def build_risks_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_header_bar(slide, prs, data["title"])

    for i, (risk, mitigation) in enumerate(data["pairs"]):
        top = 1.6 + i * 1.25
        rbox = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(5.2), Inches(0.9))
        tf = rbox.text_frame
        tf.text = f"⚠  {risk}"
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ACCENT_ORANGE

        mbox = slide.shapes.add_textbox(Inches(6.2), Inches(top), Inches(6.0), Inches(0.9))
        tf = mbox.text_frame
        tf.text = f"✓  {mitigation}"
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.color.rgb = TEAL

    add_notes(slide, data.get("notes", ""))


def build_cta_slide(prs, data):
    build_highlight_slide(prs, data)


def build_closing_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    title = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.5), Inches(1.2))
    tf = title.text_frame
    tf.text = data["title"]
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(11.5), Inches(1.0))
    tf = sub.text_frame
    tf.text = data["subtitle"]
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.color.rgb = TEAL
    p.alignment = PP_ALIGN.CENTER

    if data.get("footer"):
        foot = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.5))
        tf = foot.text_frame
        tf.text = data["footer"]
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = MID_GRAY
        p.alignment = PP_ALIGN.CENTER

    add_notes(slide, data.get("notes", ""))


BUILDERS = {
    "title": build_title_slide,
    "content": build_content_slide,
    "highlight": build_highlight_slide,
    "two_column": build_two_column_slide,
    "process": build_process_slide,
    "cards": build_cards_slide,
    "stats": build_stats_slide,
    "comparison": build_comparison_slide,
    "revenue": build_revenue_slide,
    "roadmap": build_roadmap_slide,
    "financial": build_financial_slide,
    "team": build_team_slide,
    "impact": build_impact_slide,
    "risks": build_risks_slide,
    "cta": build_cta_slide,
    "closing": build_closing_slide,
}


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_data in SLIDES:
        layout = slide_data["layout"]
        BUILDERS[layout](prs, slide_data)

    prs.save(OUTPUT)
    print(f"Created: {OUTPUT} ({len(SLIDES)} slides)")


if __name__ == "__main__":
    main()
